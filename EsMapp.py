# API key stored in Replit Secrets as GEMINI_API_KEY
import streamlit as st
import openai
import tempfile
from openai import OpenAI
import google.generativeai as genai
import edge_tts
import asyncio
import PyPDF2
from io import StringIO
import ast
from streamlit.components.v1 import html
import math
from bs4 import BeautifulSoup
import numpy as np
from difflib import SequenceMatcher
import os
import requests
import json
import fitz          # PyMuPDF  → pip install pymupdf
import base64
from PIL import Image
import io
import random

st.set_page_config(
    page_title="LleY",
    page_icon=":writing_hand:",
    layout="wide"
)
st.logo("final logo 2.png", icon_image="enlarge 1.png", size="large")


def ai_assistant(prompt, rule):
    """OpenAI-compatible wrapper for Gemini Flash"""
    try:
        client = OpenAI(
            api_key=os.environ.get("GEMINI_API_KEY", "AIzaSyCIcsMgAFh_Ki5KnIebH96zjh9WQ-M397g"),
            base_url="https://generativelanguage.googleapis.com/v1beta/openai/"
        )
        messages = []
        if rule:
            messages.append({"role": "system", "content": rule})
        if isinstance(prompt, str):
            messages.append({"role": "user", "content": prompt})
        elif isinstance(prompt, list):
            messages.extend(prompt)

        response = client.chat.completions.create(
            model="gemini-3-flash-preview",
            messages=messages,
            temperature=0.7,
            max_tokens=20000
        )
        if response.choices and response.choices[0].message.content:
            return response.choices[0].message.content
        return None
    except Exception as e:
        st.error(f"API Error: {str(e)}")
        return None


# ─────────────────────────────────────────────────────────────
#  MAIN PAGE
# ─────────────────────────────────────────────────────────────
def main_page():
    st.markdown("""
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Chakra+Petch:wght@700&display=swap');
        html, body, .stApp {
            overflow: hidden !important; height: 100vh !important;
            width: 100vw !important; margin: 0 !important;
            padding: 0 !important; background: #000 !important;
        }
        .center-container {
            position: fixed; top: 0; left: 0;
            width: 100%; height: 100%;
            display: flex; justify-content: center; align-items: center;
        }
        .main-text {
            font-family: 'Chakra Petch', sans-serif;
            font-size: clamp(3rem, 10vw, 8rem); font-weight: 700;
            color: #ff6a6a; text-transform: uppercase;
            letter-spacing: 0.5rem; text-align: center;
            position: relative; z-index: 10;
        }
        .orbit-container {
            position: absolute; top: 0; left: 0;
            width: 100%; height: 100%; pointer-events: none; z-index: 1;
        }
        .orbiting-circle {
            position: absolute; border-radius: 50%;
            background: transparent; border: 2px solid #ff69b4;
            filter: drop-shadow(0 0 5px #ff1493);
            transform-origin: center center;
        }
    </style>
    """, unsafe_allow_html=True)

    html("""
    <script>
    function waitForElm(selector) {
        return new Promise(resolve => {
            if (document.querySelector(selector)) return resolve(document.querySelector(selector));
            const observer = new MutationObserver(mutations => {
                if (document.querySelector(selector)) { observer.disconnect(); resolve(document.querySelector(selector)); }
            });
            observer.observe(document.body, { childList: true, subtree: true });
        });
    }
    waitForElm('.center-container').then((container) => {
        const orbitContainer = document.createElement('div');
        orbitContainer.className = 'orbit-container';
        container.prepend(orbitContainer);
        const centerX = window.innerWidth / 2;
        const centerY = window.innerHeight / 2;
        for (let i = 0; i < 8; i++) {
            const circle = document.createElement('div');
            circle.className = 'orbiting-circle';
            const size = 20 + (i * 15);
            const radius = 80 + (i * 60);
            const speed = 0.2 + (i * 0.05);
            const startAngle = (i * Math.PI * 2) / 8;
            circle.style.width = `${size}px`; circle.style.height = `${size}px`;
            circle.style.borderWidth = `${1 + (i * 0.3)}px`;
            let angle = startAngle;
            function animate() {
                angle += speed * 0.01;
                const x = centerX + Math.cos(angle) * radius;
                const y = centerY + Math.sin(angle) * radius;
                circle.style.left = `${x - size/2}px`; circle.style.top = `${y - size/2}px`;
                requestAnimationFrame(animate);
            }
            orbitContainer.appendChild(circle); animate();
        }
        document.body.style.overflow = 'hidden';
        window.addEventListener('scroll', () => window.scrollTo(0, 0));
    });
    </script>
    """)

    st.markdown("""
    <div class="center-container">
        <div class="main-text">LleY</div>
    </div>
    """, unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────────
#  ESSAY MAKER
# ─────────────────────────────────────────────────────────────
def esma():
    list_essay_type = [
        "Argumentative", "Persuasive", 'Explanatory', 'Descriptive',
        "Narrative", 'Cause and Effect', "Process Analysis",
        "Compare/Contrast", "Critique", "Definition", "General"
    ]
    list_level = ["Elementary", "Junior High", "Senior High", "Undergraduate", "Graduate", "Postgraduate", "PhD",
                  "Masters", "Doctorate"]
    list_speech_type = ["Casual", "Intimate", "Formal", "Frozen", "Consultative"]
    left, right = st.columns(2, vertical_alignment="top")
    languages = ["English (US)", "Filipino", "Hindi"]
    content_list_category = ["Facebook", "Youtube", "Tiktok", "Twitter/X", "Reddit", 'Instagram', "Spotify"]
    content_list_content = ["Vlog", "Blog", "Written Post", "Podcast", "Music", "Article"]
    content_list_complexity = ["Micro-Content", "Standard Content", "Premium Content", "Enterprise Content", "Legacy Content"]
    content_list_tiers = ["Beginner", "Intermediate", "Advanced", "Expert"]
    content_creator_types = [
        "The Guru", "The Everyman", "The Mad Scientist", "The Storyteller", "The Entertainer",
        "The Professor", "The Simplifier", "The Debunker", "The Coach",
        "The Hero", "The Mentor", "The Trickster", "The Villain",
        "The Reactor", "The Investigator", "The Trendsetter", "The Parodist", "The Nostalgist", "The Minimalist",
    ]
    condition_system = """You are EsMa. Strictly only write the requested essay content or content creation scripts. 
                        Do not write any other information. Meaning, only write paragraphs (unless user chose content creation). 
                        Also the output must have be clear and specific with no vague output"""

    def main():
        with st.container():
            left.subheader("EsMa by Elley")
            left.title("Your AI content generator")
        on = left.toggle("Activate Content Creation")
        if not on:
            with st.container():
                level_essay = left.selectbox("Type-Level", list_level)
                speech_type = left.selectbox("Type of Speech", list_speech_type)
                essay_type = left.selectbox("Essay Type", list_essay_type)
                speech = left.selectbox("Select Language", languages)
                word_num = left.slider("Select Number Words", min_value=0, max_value=1500, step=100)
                selected_pov = left.segmented_control('Point of View', ['First', 'Second', 'Third'], selection_mode="single")
            with st.container():
                content_prompt = left.text_area("Prompt", "", height=150)
                other_info_prompt = left.text_area("Other Instructions", "", height=70)
                if st.button("Generate Essay"):
                    if content_prompt.strip() in ("", "Generated prompt"):
                        st.warning("Please enter a valid prompt")
                    else:
                        with st.spinner("Generating your essay..."):
                            full_prompt = f"Write a comprehensive {essay_type}, point of view: {selected_pov} point of view, education level: {level_essay}, language to use:{speech}, type of speech: {speech_type}, number of minimum words: {word_num}, essay about: {content_prompt}. With extra task {other_info_prompt}"
                            essay = ai_assistant(full_prompt, condition_system)
                            if essay:
                                right.text_area("Generated Essay", value=essay, height=680)
                else:
                    right.text_area("Generated Essay", "", height=680)
        else:
            with st.container():
                content_place = left.selectbox("Where to Post?", content_list_category)
                content_type = left.selectbox("Type of Content", content_list_content)
                content_complexity = left.selectbox("Content Complexity", content_list_complexity)
                content_tier = left.selectbox("Content Tier", content_list_tiers)
                selected_character = left.selectbox("Archetype", content_creator_types)
                word_num = left.slider("Select Number Words", min_value=0, max_value=1500, step=100)
                speech = left.segmented_control('Language', languages, selection_mode="single")
            with st.container():
                content_prompt = left.text_area("Prompt", "", height=150)
                other_info_prompt = left.text_area("Other Instructions", "", height=70)
                if st.button("Generate Essay"):
                    if content_prompt.strip() in ("", "Generated prompt"):
                        st.warning("Please enter a valid prompt")
                    else:
                        with st.spinner("Generating your content..."):
                            full_prompt = f"""Write a content about: {content_prompt} in a language:{speech}, having a word length of: {word_num}, that will be posted or used in: {content_place}
having a content type: {content_type}, having a content complexity: {content_complexity}, having a content tier for: {content_tier}
and will portray a character: {selected_character}. With extra task {other_info_prompt}"""
                            content_out = ai_assistant(full_prompt, condition_system)
                            if content_out:
                                right.text_area("Generated Content", value=content_out, height=680)
                else:
                    right.text_area("Generated Content", "", height=680)

    if __name__ == "__main__":
        main()
    main()


# ─────────────────────────────────────────────────────────────
#  TTS VOICES
# ─────────────────────────────────────────────────────────────
voices_by_gender = {
    "Male": {
        "ChristopherNeural (US)": {"ShortName": "en-US-ChristopherNeural"},
        "EricNeural (US)": {"ShortName": "en-US-EricNeural"},
        "GuyNeural (US)": {"ShortName": "en-US-GuyNeural"},
        "RogerNeural (US)": {"ShortName": "en-US-RogerNeural"},
        "SteffanNeural (US)": {"ShortName": "en-US-SteffanNeural"},
        "RyanNeural (UK)": {"ShortName": "en-GB-RyanNeural"},
        "ThomasNeural (UK)": {"ShortName": "en-GB-ThomasNeural"},
        "WilliamNeural (Australian)": {"ShortName": "en-AU-WilliamNeural"},
        "MitchellNeural (New Zealand)": {"ShortName": "en-NZ-MitchellNeural"},
        "LiamNeural (Canada)": {"ShortName": "en-CA-LiamNeural"},
        "DmitryNeural (Russian)": {"ShortName": "ru-RU-DmitryNeural"},
        "HenriNeural (French)": {"ShortName": "fr-FR-HenriNeural"},
        "ConradNeural (German)": {"ShortName": "de-DE-ConradNeural"},
        "DiegoNeural (Italian)": {"ShortName": "it-IT-DiegoNeural"},
        "JorgeNeural (Mexico)": {"ShortName": "es-MX-JorgeNeural"},
        "AlvaroNeural (Spain)": {"ShortName": "es-ES-AlvaroNeural"},
        "MadhurNeural (Hindi)": {"ShortName": "hi-IN-MadhurNeural"},
        "YunjianNeural (Chinese)": {"ShortName": "zh-CN-YunjianNeural"},
        "InJoonNeural (Korean)": {"ShortName": "ko-KR-InJoonNeural"},
    },
    "Female": {
        "AriaNeural (US)": {"ShortName": "en-US-AriaNeural"},
        "JennyNeural (US)": {"ShortName": "en-US-JennyNeural"},
        "MichelleNeural (US)": {"ShortName": "en-US-MichelleNeural"},
        "AnaNeural (US)": {"ShortName": "en-US-AnaNeural"},
        "LibbyNeural (UK)": {"ShortName": "en-GB-LibbyNeural"},
        "SoniaNeural (UK)": {"ShortName": "en-GB-SoniaNeural"},
        "MaisieNeural (UK)": {"ShortName": "en-GB-MaisieNeural"},
        "NatashaNeural (Australian)": {"ShortName": "en-AU-NatashaNeural"},
        "MollyNeural (New Zealand)": {"ShortName": "en-NZ-MollyNeural"},
        "ClaraNeural (Canada)": {"ShortName": "en-CA-ClaraNeural"},
        "DeniseNeural (French)": {"ShortName": "fr-FR-DeniseNeural"},
        "KatjaNeural (German)": {"ShortName": "de-DE-KatjaNeural"},
        "ElsaNeural (Italian)": {"ShortName": "it-IT-ElsaNeural"},
        "DaliaNeural (Mexico)": {"ShortName": "es-MX-DaliaNeural"},
        "ElviraNeural (Spanish)": {"ShortName": "es-ES-ElviraNeural"},
        "SwaraNeural (Hindi)": {"ShortName": "hi-IN-SwaraNeural"},
        "XiaoxiaoNeural (Chinese)": {"ShortName": "zh-CN-XiaoxiaoNeural"},
        "SunHiNeural (Korean)": {"ShortName": "ko-KR-SunHiNeural"},
        "NanamiNeural (Japanese)": {"ShortName": "ja-JP-NanamiNeural"},
        "BlessicaNeural (Filipino)": {"ShortName": "fil-PH-BlessicaNeural"},
    },
}


async def text_to_speech(text, filename, character):
    communicate = edge_tts.Communicate(text, character)
    await communicate.save(filename)
    return filename


def tetos():
    st.subheader("TeTos by Elley")
    st.title("Free Online Text-To-Speech Tool")
    content_prompt = st.text_area("Prompt", "", height=150)
    selected_voice = st.segmented_control('Gender', ['Male', 'Female'], selection_mode="single")
    voice_final = None
    if selected_voice == "Male":
        voice_final = st.selectbox("Male Characters", list(voices_by_gender["Male"].keys()))
    elif selected_voice == "Female":
        voice_final = st.selectbox("Female Characters", list(voices_by_gender["Female"].keys()))
    if st.button("Generate Voice"):
        if content_prompt.strip() in ("", "Generated prompt"):
            st.warning("Please enter a valid prompt")
        else:
            with st.spinner("Generating your voice..."):
                filename = "voice.mp3"
                voice_arranged = voices_by_gender[selected_voice][voice_final]["ShortName"]
                asyncio.run(text_to_speech(content_prompt, filename, voice_arranged))
                with open(filename, "rb") as f:
                    audio_bytes = f.read()
                st.audio(audio_bytes, format="audio/mpeg", loop=True)


# ─────────────────────────────────────────────────────────────
#  AITO CHATBOT
# ─────────────────────────────────────────────────────────────
def aito():
    st.title("AITO")
    st.caption("AI TOol for General Purpose")
    if "messages" not in st.session_state:
        st.session_state.messages = [
            {"role": "system", "content": "You shall be named AITO a general AI TOol"}
        ]
    for msg in st.session_state.messages:
        if msg["role"] != "system":
            st.chat_message(msg["role"]).write(msg["content"])
    if prompt := st.chat_input():
        st.session_state.messages.append({"role": "user", "content": prompt})
        st.chat_message("user").write(prompt)
        response = ai_assistant(st.session_state.messages, rule=None)
        if response:
            st.session_state.messages.append({"role": "assistant", "content": response})
            st.chat_message("assistant").write(response)


# ─────────────────────────────────────────────────────────────
#  PDF HELPERS
# ─────────────────────────────────────────────────────────────
def extract_pdf_text(pdf_file: str) -> str:
    try:
        reader = PyPDF2.PdfReader(pdf_file)
        pdf_text = []
        for page in reader.pages:
            content = page.extract_text()
            if content:
                pdf_text.append(content)
        return "\n".join(pdf_text) if pdf_text else "No text could be extracted from the PDF."
    except Exception as f:
        st.error(f"Error reading PDF: {str(f)}")
        return None


def extract_pptx_text(pptx_file):
    try:
        from pptx import Presentation
        prs = Presentation(pptx_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text) if text else "No text could be extracted from the PPTX."
    except ImportError:
        st.error("Please install python-pptx package: pip install python-pptx")
        return None
    except Exception as h:
        st.error(f"Error reading PPTX: {str(h)}")
        return None


# ─────────────────────────────────────────────────────────────
#  PDF TO QUIZ
# ─────────────────────────────────────────────────────────────
def pdf2quiz():
    system_condition_mcq = """You are a system only for creating multiple choice quiz python dictionaries. 
                        Return ONLY a properly formatted Python dictionary with no additional text or explanation.
                        Format:
                        {
                            "1": {
                                "question": "Question text",
                                "a": "Option A",
                                "b": "Option B",
                                "c": "Option C",
                                "d": "Option D",
                                "answer_key": "correct_letter"
                            },
                            ...
                        }"""
    system_condition_open = """You are a system only for creating open-ended quiz python dictionaries. 
                        Return ONLY a properly formatted Python dictionary with no additional text or explanation.
                        Format:
                        {
                            "1": {
                                "question": "Question text",
                                "type": "definition/enumeration/essay",
                                "model_answer": "The ideal answer that would score 10/10 (3-5 sentences minimum)",
                                "scoring_criteria": ["Key point 1", "Key point 2", "Key point 3"]
                            },
                            ...
                        }"""
    scoring_system = """You are an expert grader. Return ONLY this format:
    {
        "score": [1-10],
        "explanation": "Paragraph analyzing strengths/weaknesses",
        "feedback": "3 specific improvement suggestions",
        "key_matches": ["list", "of", "matched", "concepts"],
        "missing_points": ["list", "of", "missing", "elements"]
    }"""

    if 'quiz' not in st.session_state:
        st.session_state.quiz = {
            'data': None, 'answers': {}, 'submitted': False,
            'file_processed': None, 'file_type': None,
            'quiz_type': 'multiple_choice', 'scores': {}, 'custom_topic': None
        }

    st.title("📝 Smart Quiz Generator")
    st.subheader("Create quizzes from files or custom topics")

    quiz_type = st.radio("Quiz Type", options=["Multiple Choice", "Open-Ended"], key="quiz_type_selector")
    st.session_state.quiz['quiz_type'] = 'multiple_choice' if quiz_type == "Multiple Choice" else 'open_ended'

    content_source = st.radio("Content Source", options=["Upload File", "Custom Topic"], horizontal=True, key="content_source")

    if content_source == "Upload File":
        uploaded_file = st.file_uploader("Upload PDF or PPTX", type=["pdf", "pptx"])
        custom_topic = None
    else:
        custom_topic = st.text_area("Enter your custom topic/subject:", placeholder="e.g., Machine Learning...", key="custom_topic")
        uploaded_file = None
        st.session_state.quiz['custom_topic'] = custom_topic

    number_quiz = st.number_input("Number of questions", min_value=1, max_value=100, value=6)
    content_prompt = st.text_area("Extra Prompt", "", height=90)

    if st.button("Generate Quiz", type="primary"):
        if (uploaded_file is None) and (custom_topic is None or custom_topic.strip() == ""):
            st.warning("Please either upload a file or enter a custom topic")
        else:
            with st.spinner(f"Generating {quiz_type.lower()} quiz..."):
                extracted_text = ""
                if uploaded_file:
                    with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
                        tmp_file.write(uploaded_file.getvalue())
                        tmp_file_path = tmp_file.name
                    extracted_text = extract_pdf_text(tmp_file_path) if uploaded_file.name.endswith('.pdf') else extract_pptx_text(tmp_file_path)

                full_prompt = f"Create {number_quiz} {'multiple choice' if st.session_state.quiz['quiz_type'] == 'multiple_choice' else 'open-ended'} questions. Additional info: {content_prompt}"
                if custom_topic:
                    full_prompt += f" about: {custom_topic}"
                if extracted_text and not extracted_text.startswith("No text"):
                    full_prompt += f"\nText content: {extracted_text[:10000]}"

                system_prompt = system_condition_mcq if st.session_state.quiz['quiz_type'] == 'multiple_choice' else system_condition_open
                content_out = ai_assistant(full_prompt, system_prompt)

                if content_out:
                    try:
                        clean_output = content_out.strip()
                        if clean_output.startswith("```"):
                            clean_output = clean_output.split("```")[1]
                            if clean_output.startswith(("python", "json")):
                                clean_output = clean_output[6:].strip()
                        clean_output = clean_output.strip().strip('"').strip("'")
                        if not clean_output.endswith('}'):
                            clean_output += '}'
                        quiz_data = ast.literal_eval(clean_output)
                        st.session_state.quiz.update({
                            'data': quiz_data,
                            'answers': {q_num: None for q_num in quiz_data},
                            'file_processed': uploaded_file,
                            'file_type': "PDF" if uploaded_file and uploaded_file.name.endswith('.pdf') else "PPTX" if uploaded_file else "Custom Topic",
                            'scores': {}, 'submitted': False
                        })
                        st.rerun()
                    except Exception as e:
                        st.error(f"Error processing quiz: {str(e)}")
                        st.code(content_out)

    if st.session_state.quiz['data']:
        st.divider()
        st.subheader(f"{quiz_type} Quiz")
        all_answered = True
        for q_num, question in st.session_state.quiz['data'].items():
            st.markdown(f"### Question {q_num}")
            st.write(question['question'])
            if st.session_state.quiz['quiz_type'] == 'multiple_choice':
                options = [question['a'], question['b'], question['c'], question['d']]
                current_answer = st.session_state.quiz['answers'].get(q_num)
                user_choice = st.radio("Select your answer:", options, key=f"q_{q_num}_mcq",
                                       index=options.index(current_answer) if current_answer in options else None)
                if user_choice != current_answer:
                    st.session_state.quiz['answers'][q_num] = user_choice
                    st.rerun()
            else:
                current_answer = st.session_state.quiz['answers'].get(q_num, "")
                user_answer = st.text_area("Your answer:", value=current_answer, key=f"q_{q_num}_open", height=150)
                if user_answer != current_answer:
                    st.session_state.quiz['answers'][q_num] = user_answer
                    st.rerun()
            if st.session_state.quiz['answers'].get(q_num) is None:
                all_answered = False
            if st.session_state.quiz['submitted']:
                if st.session_state.quiz['quiz_type'] == 'multiple_choice':
                    correct_answer = question[question['answer_key']]
                    if st.session_state.quiz['answers'][q_num] == correct_answer:
                        st.success("✓ Correct!")
                    else:
                        st.error(f"✗ Incorrect. The correct answer is: {correct_answer}")
                else:
                    score_data = st.session_state.quiz['scores'].get(q_num, {})
                    st.markdown(f"**Score:** {score_data.get('score', 0)}/10")
                    with st.expander("📝 Detailed Feedback"):
                        st.write(score_data.get('explanation', 'No explanation available'))
                    with st.expander("📚 Model Answer"):
                        st.write(question.get('model_answer', 'Not available'))
            st.markdown("---")

        col1, col2 = st.columns(2)
        with col1:
            if not st.session_state.quiz['submitted'] and st.button("Submit Answers", disabled=not all_answered):
                if st.session_state.quiz['quiz_type'] == 'open_ended':
                    with st.spinner("Evaluating answers..."):
                        scores = {}
                        for q_num, question in st.session_state.quiz['data'].items():
                            user_answer = st.session_state.quiz['answers'][q_num]
                            if not user_answer or len(user_answer.strip()) < 10:
                                scores[q_num] = {"score": 1, "explanation": "Answer too short", "feedback": "Provide more detail", "key_matches": [], "missing_points": []}
                                continue
                            prompt = f"Question: {question['question']}\nModel: {question.get('model_answer','')}\nStudent: {user_answer}"
                            try:
                                score_data = ai_assistant(prompt, scoring_system)
                                clean_data = score_data.strip().strip('```').strip()
                                scores[q_num] = json.loads(clean_data)
                            except:
                                scores[q_num] = {"score": 5, "explanation": "Evaluation failed", "feedback": "Compare with model answer", "key_matches": [], "missing_points": []}
                        st.session_state.quiz['scores'] = scores
                st.session_state.quiz['submitted'] = True
                st.rerun()
        with col2:
            if st.button("Reset Quiz"):
                st.session_state.quiz = {'data': None, 'answers': {}, 'submitted': False, 'file_processed': None, 'file_type': None, 'quiz_type': 'multiple_choice', 'scores': {}, 'custom_topic': None}
                st.rerun()

        if st.session_state.quiz['submitted']:
            st.divider()
            if st.session_state.quiz['quiz_type'] == 'multiple_choice':
                correct = sum(1 for q_num, question in st.session_state.quiz['data'].items()
                              if st.session_state.quiz['answers'][q_num] == question[question['answer_key']])
                st.success(f"🎯 Your score: {correct}/{len(st.session_state.quiz['data'])} ({correct/len(st.session_state.quiz['data'])*100:.1f}%)")
            else:
                total = sum(score['score'] for score in st.session_state.quiz['scores'].values())
                max_score = 10 * len(st.session_state.quiz['data'])
                st.success(f"🎯 Total Score: {total}/{max_score} ({total/max_score*100:.1f}%)")


# ─────────────────────────────────────────────────────────────
#  ORIGINALITY CHECKER
# ─────────────────────────────────────────────────────────────
def extract_text_from_file(uploaded_file):
    text = ""
    try:
        if uploaded_file.name.endswith('.pdf'):
            reader = PyPDF2.PdfReader(uploaded_file)
            text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif uploaded_file.name.endswith('.txt'):
            text = uploaded_file.read().decode('utf-8')
    except Exception as e:
        st.error(f"Error extracting text: {str(e)}")
    return text


def turnitin_knockoff():
    st.title("🔍 Originality Checker")
    st.caption("Academic integrity analysis inspired by Turnitin")
    input_method = st.radio("Input Method", ["Text Input", "File Upload", "Website URL"], horizontal=True)
    text = ""
    if input_method == "Text Input":
        text = st.text_area("Enter text to analyze", height=200)
    elif input_method == "File Upload":
        uploaded_file = st.file_uploader("Upload document", type=["pdf", "docx", "pptx", "txt"])
        if uploaded_file:
            with st.spinner("Extracting text..."):
                text = extract_text_from_file(uploaded_file)
    else:
        url = st.text_input("Enter website URL")
        if url:
            with st.spinner("Fetching website..."):
                try:
                    response = requests.get(url)
                    soup = BeautifulSoup(response.text, 'html.parser')
                    text = soup.get_text()
                except Exception as e:
                    st.error(f"Error fetching URL: {str(e)}")

    if not text.strip():
        st.warning("No text found to analyze")
        return

    st.metric("Word Count", len(text.split()))
    analysis_type = st.radio("Analysis Mode", ["Quick Check", "Deep Analysis"], horizontal=True)
    analysis_depth = 3000 if analysis_type == "Quick Check" else 10000

    if st.button("Run Originality Check"):
        with st.spinner("Analyzing content..."):
            st.subheader("🤖 AI Detection Score")
            ai_prompt = f"""Analyze this text for AI-generated patterns: {text[:analysis_depth]}
Return ONLY a JSON object: {{"score": 0-100, "flagged_passages": [["phrase", score], ...], "explanation": "string"}}"""
            ai_result = ai_assistant(ai_prompt, "You are an AI content detector. Return ONLY valid JSON.")
            if ai_result:
                try:
                    clean = ai_result.strip().lstrip("```json").lstrip("```").rstrip("```")
                    ai_data = json.loads(clean)
                    st.progress(ai_data["score"] / 100)
                    st.metric("AI Likelihood Score", f"{ai_data['score']}%")
                    with st.expander("AI Analysis Details"):
                        st.caption(ai_data.get("explanation", ""))
                except Exception as e:
                    st.error(f"Error: {str(e)}")

            st.subheader("📝 Internal Similarity")
            sim_prompt = f"""Analyze repetition in: {text[:analysis_depth]}
Return ONLY JSON: {{"repetition_score": 0-100, "most_repeated_phrases": [{{"phrase": "...", "count": N}}], "suggestions": "string"}}"""
            sim_result = ai_assistant(sim_prompt, "You analyze text repetition. Return ONLY valid JSON.")
            if sim_result:
                try:
                    clean = sim_result.strip().lstrip("```json").lstrip("```").rstrip("```")
                    sim_data = json.loads(clean)
                    st.progress(sim_data["repetition_score"] / 100)
                    st.metric("Repetition Score", f"{sim_data['repetition_score']}%")
                    st.info("Suggestions: " + sim_data.get("suggestions", ""))
                except Exception as e:
                    st.error(f"Error: {str(e)}")


# ─────────────────────────────────────────────────────────────
#  🔣 SYMBOL QUIZ  (image → type the answer)
# ─────────────────────────────────────────────────────────────
def symbol_quiz():
    st.title("🔣 Symbol Quiz")
    st.caption("See a symbol image — type what it means. Upload your symbol PDF to begin.")

    # Session state
    if "sq" not in st.session_state:
        st.session_state.sq = {
            "pairs":        [],
            "quiz":         [],
            "user_answers": {},
            "submitted":    False,
            "num_q":        10,
        }
    sq = st.session_state.sq

    # Extract embedded images from PDF and pair with nearest text label
    def extract_image_pairs(uploaded_file) -> list:
        raw = uploaded_file.read()
        doc = fitz.open(stream=raw, filetype="pdf")
        pairs = []
        seen_xrefs = set()

        for page in doc:
            page_area   = page.rect.width * page.rect.height
            img_list    = page.get_images(full=True)
            text_blocks = [b for b in page.get_text("blocks") if b[6] == 0 and b[4].strip()]

            for img_info in img_list:
                xref = img_info[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)

                base_img  = doc.extract_image(xref)
                img_bytes = base_img["image"]
                iw        = base_img.get("width", 0)
                ih        = base_img.get("height", 0)

                if iw < 15 or ih < 15:
                    continue
                if iw * ih > 0.45 * page_area:
                    continue

                rects = page.get_image_rects(xref)
                if not rects:
                    continue
                img_rect = rects[0]
                icx = (img_rect.x0 + img_rect.x1) / 2
                icy = (img_rect.y0 + img_rect.y1) / 2

                if not text_blocks:
                    continue

                # Layout: [label] [image] [image] [label]
                # The label is to the LEFT of the image, on the same row.
                # Step 1: filter to blocks that are:
                #   - horizontally to the left of the image (right edge of text <= left edge of image + small overlap)
                #   - vertically overlapping the image row (centres within image height + tolerance)
                img_h      = img_rect.y1 - img_rect.y0
                row_top    = img_rect.y0 - img_h * 0.6
                row_bottom = img_rect.y1 + img_h * 0.6

                left_candidates = [
                    b for b in text_blocks
                    if b[2] <= img_rect.x0 + 20           # text right edge is left of image
                    and row_top <= (b[1] + b[3]) / 2 <= row_bottom  # same row
                ]

                # Step 2: if nothing to the left, fall back to closest block overall
                # (handles edge cases like single-column pages)
                if left_candidates:
                    # Pick the one whose right edge is closest to the image left edge
                    best = min(left_candidates, key=lambda b: img_rect.x0 - b[2])
                else:
                    def score(b):
                        bx0, by0, bx1, by1 = b[0], b[1], b[2], b[3]
                        vert  = min(abs(by0 - img_rect.y1), abs(by1 - img_rect.y0))
                        horiz = abs((bx0 + bx1) / 2 - icx) * 0.3
                        return vert + horiz
                    best = min(text_blocks, key=score)
                    if score(best) > 120:
                        continue

                label = best[4].strip().replace("\n", " ")
                if not label or len(label) > 120:
                    continue
                if sum(c.isalpha() for c in label) < 2:
                    continue

                pairs.append({"img_bytes": img_bytes, "answer": label})

        doc.close()
        seen, unique = set(), []
        for p in pairs:
            key = p["answer"].lower().strip()
            if key not in seen:
                seen.add(key)
                unique.append(p)
        return unique

    # ── Render image from bytes ───────────────────────────────
    def show_image(img_bytes: bytes, size: int = 180):
        try:
            img = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
            img.thumbnail((size, size), Image.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            b64 = base64.b64encode(buf.getvalue()).decode()
            st.markdown(
                f"""
                <div style="
                    display:flex; justify-content:center; align-items:center;
                    padding:16px; background:#1a1a2e; border-radius:12px;
                    border:2px solid #ff6a6a; margin-bottom:8px;">
                  <img src="data:image/png;base64,{b64}"
                       style="max-width:{size}px; max-height:{size}px;
                              image-rendering:pixelated;" />
                </div>
                """,
                unsafe_allow_html=True,
            )
        except Exception as e:
            st.error(f"Could not render image: {e}")

    # ── Fuzzy match helper ────────────────────────────────────
    def is_correct(user: str, correct: str) -> tuple[bool, bool, int]:
        """Returns (exact_match, close_match, pct)"""
        u = user.lower().strip()
        c = correct.lower().strip()
        pct = int(SequenceMatcher(None, u, c).ratio() * 100)
        return u == c, pct >= 70, pct

    # ─────────────────────────────────────────────────────────
    # SIDEBAR — upload + settings
    # ─────────────────────────────────────────────────────────
    with st.sidebar:
        st.header("⚙️ Symbol Quiz Settings")
        uploaded_pdf  = st.file_uploader("📄 Upload Symbol PDF", type=["pdf"])
        num_questions = st.slider("Questions per round", min_value=3, max_value=40, value=10)

        if uploaded_pdf and st.button("📥 Load PDF & Build Quiz", type="primary"):
            with st.spinner("Extracting symbol images from PDF…"):
                pairs = extract_image_pairs(uploaded_pdf)

            if not pairs:
                st.error("No symbol images found. Make sure the PDF has embedded symbol drawings with text labels nearby.")
            else:
                st.success(f"✅ Found {len(pairs)} symbol pairs!")
                pool = pairs.copy()
                random.shuffle(pool)
                sq["pairs"]        = pairs
                sq["quiz"]         = pool[:min(num_questions, len(pool))]
                sq["user_answers"] = {}
                sq["submitted"]    = False
                sq["num_q"]        = num_questions
                st.rerun()

        if sq["pairs"]:
            st.caption(f"Library: {len(sq['pairs'])} symbols loaded")
            if st.button("🔀 New Random Round"):
                pool = sq["pairs"].copy()
                random.shuffle(pool)
                sq["quiz"]         = pool[:min(num_questions, len(pool))]
                sq["user_answers"] = {}
                sq["submitted"]    = False
                st.rerun()

    # ─────────────────────────────────────────────────────────
    # EMPTY STATE
    # ─────────────────────────────────────────────────────────
    if not sq["quiz"]:
        st.info("👈 Upload your symbol PDF in the sidebar, then click **Load PDF & Build Quiz**.")
        st.markdown("""
        **How it works:**
        1. Upload a PDF where each symbol is an **embedded image** next to its label/name
        2. The app automatically pairs each image with its nearest text label
        3. You see the symbol — type the answer — get instant feedback after submitting

        > 💡 Works best with PDFs that have clear image + label layouts  
        > (e.g. electrical symbols, chemical hazard symbols, music notation, map legends)
        """)
        return

    # ─────────────────────────────────────────────────────────
    # SCORE BANNER (shown after submit)
    # ─────────────────────────────────────────────────────────
    total_q = len(sq["quiz"])

    if sq["submitted"]:
        results = [
            is_correct(sq["user_answers"].get(i, ""), item["answer"])
            for i, item in enumerate(sq["quiz"])
        ]
        score = sum(1 for exact, close, _ in results if exact or close)
        pct   = score / total_q * 100
        emoji = "🏆" if pct >= 90 else "🎉" if pct >= 75 else "👍" if pct >= 55 else "📚"

        st.markdown(
            f"""
            <div style="background:linear-gradient(135deg,#1a1a2e,#16213e);
                        border:2px solid #ff6a6a; border-radius:14px;
                        padding:20px; text-align:center; margin-bottom:24px;">
              <div style="font-size:2.5rem;">{emoji}</div>
              <div style="font-size:1.8rem; font-weight:700; color:#ff6a6a;">
                {score} / {total_q}
              </div>
              <div style="color:#aaa; font-size:1rem;">{pct:.1f}% correct</div>
            </div>
            """,
            unsafe_allow_html=True,
        )

    else:
        answered = sum(1 for v in sq["user_answers"].values() if str(v).strip())
        st.markdown(
            f'<p style="color:#aaa; margin-bottom:8px;">Progress: '
            f'<b style="color:#ff6a6a">{answered}/{total_q}</b> answered</p>',
            unsafe_allow_html=True,
        )

    st.divider()

    # ─────────────────────────────────────────────────────────
    # QUIZ CARDS
    # ─────────────────────────────────────────────────────────
    all_answered = True

    for i, item in enumerate(sq["quiz"]):
        with st.container(border=True):
            left_col, right_col = st.columns([1, 2], gap="large")

            with left_col:
                st.markdown(f"**#{i+1}**")
                show_image(item["img_bytes"])

            with right_col:
                st.markdown("**What does this symbol represent?**")

                cur = sq["user_answers"].get(i, "")

                # After submit: show read-only answer + feedback
                if sq["submitted"]:
                    exact, close, pct = is_correct(cur, item["answer"])
                    st.text_input(
                        "Your answer",
                        value=cur,
                        key=f"sq_done_{i}",
                        disabled=True,
                        label_visibility="collapsed",
                    )
                    if exact:
                        st.success(f"✅ **Correct!** — {item['answer']}")
                    elif close:
                        st.warning(f"🟡 **Close enough ({pct}%)** — correct: **{item['answer']}**")
                    else:
                        st.error(f"❌ **Wrong** — correct answer: **{item['answer']}**")

                # Before submit: live input
                else:
                    ans = st.text_input(
                        "Your answer",
                        value=cur,
                        key=f"sq_q_{i}",
                        placeholder="Type the meaning or name…",
                        label_visibility="collapsed",
                    )
                    if ans != cur:
                        sq["user_answers"][i] = ans
                        # no rerun on every keystroke — let Streamlit handle naturally

                    if not str(sq["user_answers"].get(i, "")).strip():
                        all_answered = False

    # ─────────────────────────────────────────────────────────
    # ACTION BUTTONS
    # ─────────────────────────────────────────────────────────
    st.divider()
    btn_col1, btn_col2 = st.columns(2)

    with btn_col1:
        if not sq["submitted"]:
            if st.button(
                "✅ Submit All Answers",
                type="primary",
                use_container_width=True,
                disabled=not all_answered,
            ):
                sq["submitted"] = True
                st.rerun()

    with btn_col2:
        label = "🔄 Try Again (same set)" if sq["submitted"] else "🔄 Clear Answers"
        if st.button(label, use_container_width=True):
            sq["user_answers"] = {}
            sq["submitted"]    = False
            st.rerun()


# ─────────────────────────────────────────────────────────────
#  ABOUT
# ─────────────────────────────────────────────────────────────
def about():
    st.title("About LleY")
    st.markdown("""
    **LleY** is a suite of AI-powered tools for students and creators.

    | Tool | Description |
    |------|-------------|
    | AITO | General-purpose AI chatbot |
    | Essay Maker | AI essay and content generator |
    | Text To Speech | 200+ voices via Edge TTS |
    | PDF To Quiz | Generate quizzes from documents |
    | Originality Checker | AI + plagiarism detection |
    | **Symbol Quiz** | **Upload a symbol PDF and quiz yourself!** |
    """)


# ─────────────────────────────────────────────────────────────
#  TRANS MAKER (PPTX / PDF / DOCX → condensed "trans" PDF)
# ─────────────────────────────────────────────────────────────
def trans_maker():
    import textwrap
    from fpdf import FPDF

    st.title("📚 Trans Maker")
    st.caption("Convert PPTX, PDF, or DOCX lecture files into a condensed, multi-column 'trans' PDF — with images included.")

    trans_system = """You are "Trans Maker", an academic note-condenser used by college students
to create "trans" (transcription/reviewer notes) from lecture slides or documents.

STRICT RULES:
1. Base your output ONLY on the given source text. Do NOT add outside facts, examples,
   definitions, or information that is not present in the source.
2. If the source text is empty, sparse, or just a title/header, keep the output minimal —
   do not invent content to fill space.
3. Write in clear, concise, easy-to-understand bullet points suitable for reviewing/studying.
4. Use short phrases, not long paragraphs. Preserve key terms, definitions, formulas,
   and numbers exactly as given in the source.
5. Do NOT include meta-commentary, labels like "Summary:", or any text besides the JSON.

Return ONLY valid JSON in this exact format, nothing else:
{"title": "short title for this section/slide", "bullets": ["point 1", "point 2", ...]}"""

    # ── extraction helpers ──────────────────────────────────
    def extract_units_pdf(file_bytes):
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        units = []
        for page in doc:
            text = page.get_text().strip()
            images = []
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                try:
                    base_img = doc.extract_image(xref)
                    if base_img["width"] >= 40 and base_img["height"] >= 40:
                        images.append(base_img["image"])
                except Exception:
                    pass
            if text or images:
                units.append({"text": text, "images": images})
        doc.close()
        return units

    def extract_units_pptx(file_bytes):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(io.BytesIO(file_bytes))
        units = []
        for slide in prs.slides:
            text_parts, images = [], []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in para.runs)
                        if t.strip():
                            text_parts.append(t.strip())
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        images.append(shape.image.blob)
                    except Exception:
                        pass
            text = "\n".join(text_parts).strip()
            if text or images:
                units.append({"text": text, "images": images})
        return units

    def extract_units_docx(file_bytes):
        try:
            import docx
        except ImportError:
            st.error("Please install python-docx: pip install python-docx")
            return []
        d = docx.Document(io.BytesIO(file_bytes))
        images = []
        for rel in d.part.rels.values():
            if "image" in rel.reltype:
                try:
                    images.append(rel.target_part.blob)
                except Exception:
                    pass
        full_text = "\n".join(p.text.strip() for p in d.paragraphs if p.text.strip())
        chunk_size = 1400
        chunks = [full_text[i:i + chunk_size] for i in range(0, len(full_text), chunk_size)] or [""]
        units = [{"text": c, "images": []} for c in chunks]
        if images:
            per = max(1, len(images) // len(units))
            for idx, img in enumerate(images):
                units[min(idx // per, len(units) - 1)]["images"].append(img)
        return [u for u in units if u["text"] or u["images"]]

    # ── AI condensation ──────────────────────────────────────
    def condense_unit(text):
        if not text.strip():
            return {"title": "", "bullets": []}
        prompt = f"Source content:\n{text[:6000]}"
        out = ai_assistant(prompt, trans_system)
        if not out:
            return {"title": "", "bullets": [text[:300]]}
        clean = out.strip().strip("`")
        if clean.lower().startswith("json"):
            clean = clean[4:].strip()
        try:
            data = json.loads(clean)
            title = str(data.get("title", "")).strip()
            bullets = [str(b).strip() for b in data.get("bullets", []) if str(b).strip()]
            return {"title": title, "bullets": bullets}
        except Exception:
            return {"title": "", "bullets": [clean[:400]]}

    # ── PDF helper: sanitize for core (latin-1) fonts ───────
    def safe(t):
        return t.encode("latin-1", "replace").decode("latin-1")

    # ── multi-column PDF builder ────────────────────────────
    class ColumnPDF(FPDF):
        def __init__(self, num_cols=2):
            super().__init__(format="A4")
            self.num_cols = num_cols
            self.margin_l = 10
            self.gutter = 6
            self.set_auto_page_break(False)
            self.set_margins(self.margin_l, 10, self.margin_l)
            self.col_width = (210 - 2 * self.margin_l - (num_cols - 1) * self.gutter) / num_cols
            self.col_idx = 0
            self.add_page()
            self.col_y = [12] * num_cols

        def col_x(self, idx):
            return self.margin_l + idx * (self.col_width + self.gutter)

        def next_column(self):
            self.col_idx += 1
            if self.col_idx >= self.num_cols:
                self.col_idx = 0
                self.add_page()
                self.col_y = [12] * self.num_cols

        def add_section(self, title, bullets, images):
            chars_per_line = max(10, int(self.col_width / 1.7))
            line_h = 4.6
            est_lines = len(textwrap.wrap(title, chars_per_line)) or 1
            for b in bullets:
                est_lines += len(textwrap.wrap("• " + b, chars_per_line)) or 1
            text_h = est_lines * line_h + 4

            img_heights, tmp_paths = [], []
            for img_bytes in images:
                try:
                    im = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    w, h = im.size
                    ih = min(self.col_width * h / w, 75)
                    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                    im.save(tmp.name, format="PNG")
                    tmp_paths.append(tmp.name)
                    img_heights.append(ih + 3)
                except Exception:
                    img_heights.append(0)
                    tmp_paths.append(None)

            total_h = text_h + sum(img_heights)
            page_h = 297 - 10

            if self.col_y[self.col_idx] + min(total_h, page_h - 14) > page_h:
                self.next_column()

            x = self.col_x(self.col_idx)
            y = self.col_y[self.col_idx]
            self.set_xy(x, y)

            if title:
                self.set_font("Helvetica", "B", 10)
                self.set_x(x)
                self.multi_cell(self.col_width, line_h, safe(title))

            self.set_font("Helvetica", "", 8.5)
            for b in bullets:
                self.set_x(x)
                self.multi_cell(self.col_width, line_h, safe("• " + b))

            cur_y = self.get_y()
            for tmp_path, ih in zip(tmp_paths, img_heights):
                if not tmp_path or ih == 0:
                    continue
                if cur_y + ih > page_h:
                    self.col_y[self.col_idx] = cur_y + 2
                    self.next_column()
                    x = self.col_x(self.col_idx)
                    cur_y = self.col_y[self.col_idx]
                    self.set_xy(x, cur_y)
                self.image(tmp_path, x=x, y=cur_y, w=self.col_width)
                cur_y += ih
                os.remove(tmp_path)

            self.col_y[self.col_idx] = cur_y + 4

    # ── UI ───────────────────────────────────────────────────
    uploaded_file = st.file_uploader("Upload lecture file", type=["pdf", "pptx", "docx"])
    num_cols = st.select_slider("Columns per page (more = fits more, smaller text)", options=[2, 3], value=2)
    title_input = st.text_input("Trans Title (optional)", "")

    if st.button("Generate Trans", type="primary"):
        if not uploaded_file:
            st.warning("Please upload a PDF, PPTX, or DOCX file")
            return
        file_bytes = uploaded_file.getvalue()
        ext = uploaded_file.name.split(".")[-1].lower()

        with st.spinner("Extracting content from file..."):
            if ext == "pdf":
                units = extract_units_pdf(file_bytes)
            elif ext == "pptx":
                units = extract_units_pptx(file_bytes)
            elif ext == "docx":
                units = extract_units_docx(file_bytes)
            else:
                st.error("Unsupported file type")
                return

        if not units:
            st.error("No readable content found in the uploaded file")
            return

        progress = st.progress(0, text="Condensing content into trans notes...")
        sections = []
        for i, unit in enumerate(units):
            condensed = condense_unit(unit["text"])
            condensed["images"] = unit["images"]
            sections.append(condensed)
            progress.progress((i + 1) / len(units), text=f"Processing {i + 1}/{len(units)}...")
        progress.empty()

        with st.spinner("Building trans PDF..."):
            pdf = ColumnPDF(num_cols=num_cols)
            if title_input.strip():
                pdf.set_font("Helvetica", "B", 14)
                pdf.set_xy(10, 10)
                pdf.cell(0, 8, safe(title_input.strip()), ln=1)
                pdf.col_y = [pdf.get_y() + 2] * num_cols

            for sec in sections:
                if not sec["title"] and not sec["bullets"] and not sec.get("images"):
                    continue
                pdf.add_section(sec["title"], sec["bullets"], sec.get("images", []))

            pdf_bytes = bytes(pdf.output())

        st.success(f"✅ Trans generated from {len(sections)} section(s)!")
        st.download_button(
            "📥 Download Trans PDF",
            data=pdf_bytes,
            file_name=f"{(title_input.strip() or uploaded_file.name.rsplit('.', 1)[0])}_trans.pdf",
            mime="application/pdf"
        )


# ─────────────────────────────────────────────────────────────
#  NAVIGATION
# ─────────────────────────────────────────────────────────────
pages = {
    "Tools": [
        st.Page(main_page,          title="Home"),
        st.Page(aito,               title="AITO"),
        st.Page(esma,               title="Essay Maker"),
        st.Page(tetos,              title="Text To Speech"),
        st.Page(pdf2quiz,           title="Pdf To Quiz"),
        st.Page(turnitin_knockoff,  title="Originality Checker"),
        st.Page(symbol_quiz,        title="Symbol Quiz"),
        st.Page(trans_maker,        title="Transes Generator")# ← NEW
    ],
    "About": [st.Page(about, title="About")],
}

pg = st.navigation(pages)
pg.run()

